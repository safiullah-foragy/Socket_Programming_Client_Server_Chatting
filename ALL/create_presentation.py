from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(236, 72, 153)
BLUE = RGBColor(59, 130, 246)
GREEN = RGBColor(34, 197, 94)
ORANGE = RGBColor(249, 115, 22)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(30, 30, 30)
GRAY = RGBColor(100, 100, 100)

def add_title_slide():
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient-like background (dark)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 40)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Student Result Management &\nUniversity Information System"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PURPLE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Web-Based Client-Server Chat Application"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.italic = True
    subtitle_frame.paragraphs[0].font.color.rgb = PINK
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Author info
    info_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1.5))
    info_frame = info_box.text_frame
    info_text = info_frame.text = """Presented by: MD. SAFIULLAH FARAJY (2102051)
Course: Client-Server Programming
Institution: Patuakhali Science and Technology University
Date: November 22, 2025"""
    info_frame.paragraphs[0].font.size = Pt(14)
    info_frame.paragraphs[0].font.color.rgb = WHITE
    info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Key highlights
    highlights_box = slide.shapes.add_textbox(Inches(2), Inches(6.3), Inches(6), Inches(0.8))
    highlights_frame = highlights_box.text_frame
    highlights_frame.text = "💬 Interactive Chat  |  🎓 84+ Students  |  🏛️ 24 Universities  |  🎨 Stylish UI"
    highlights_frame.paragraphs[0].font.size = Pt(14)
    highlights_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 255)
    highlights_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_overview_slide():
    """Slide 2: Project Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    
    # Title
    title = slide.shapes.title
    title.text = "Project Overview & Objectives"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PURPLE
    
    # Content
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = "What is this Project?"
    
    p = tf.add_paragraph()
    p.text = "A web-based client-server application enabling conversational interaction for:"
    p.level = 0
    p.font.size = Pt(16)
    
    # Student Management
    p = tf.add_paragraph()
    p.text = "1. Student Result Management System"
    p.level = 0
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    for item in ["Store and retrieve student academic records", 
                 "Search by Student ID or Name",
                 "CGPA tracking (2.50 - 4.00 scale)",
                 "Database with 84 student records"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(14)
    
    # University Portal
    p = tf.add_paragraph()
    p.text = "2. University Information Portal"
    p.level = 0
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    for item in ["24 public universities database",
                 "Search by name or short code",
                 "Complete address and details"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(14)
    
    # Tech stack
    p = tf.add_paragraph()
    p.text = "Technology Stack: Python Flask, SQLite3, HTML5, CSS3, JavaScript"
    p.level = 0
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GRAY

def add_socket_programming_slide():
    """Slide 3: Socket Programming Components"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Socket Programming Components"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PURPLE
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = "Network Communication Implementation"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    # Socket components
    p = tf.add_paragraph()
    p.text = "🔌 HTTP Protocol (Application Layer)"
    p.font.size = Pt(18)
    p.font.color.rgb = GREEN
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Uses HTTP/1.1 protocol over TCP sockets for reliable communication"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "POST method for sending chat messages to /chat endpoint"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "⚙️ TCP/IP Stack (Transport Layer)"
    p.font.size = Pt(18)
    p.font.color.rgb = BLUE
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Flask server binds to TCP socket on 0.0.0.0:5000"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Reliable, connection-oriented stream sockets (SOCK_STREAM)"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Three-way handshake for connection establishment"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "📡 Client-Server Socket Communication"
    p.font.size = Pt(18)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Client: Browser creates socket connection via AJAX/Fetch API"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Server: Flask handles multiple concurrent connections via WSGI"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "JSON serialization for data exchange over socket"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "🔒 Socket Features: CORS enabled, persistent connections, error handling"
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = PINK

def add_architecture_slide():
    """Slide 4: System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PURPLE
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = "Client-Server Architecture Design"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    # Architecture layers
    p = tf.add_paragraph()
    p.text = "🖥️ Client Side (Web Browser)"
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "JavaScript creates socket connection using Fetch API"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Sends HTTP POST requests over TCP socket to port 5000"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "⚙️ Server Side (Flask Application)"
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "Flask binds to socket: 0.0.0.0:5000 (all network interfaces)"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Accepts incoming socket connections from multiple clients"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Routes HTTP requests to appropriate handlers"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "💾 Data Layer"
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "SQLite Database with 84 student records"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Universities Module with 24 institutions"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Socket Flow: Client Socket → TCP Connection → Server Socket → Response"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = PINK

def add_database_slide():
    """Slide 5: Database Design"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Database Design & Data Management"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PURPLE
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = "Results Table Structure"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    # Schema
    for field in ["id INTEGER PRIMARY KEY - Student ID (unique)",
                  "name TEXT NOT NULL - Full Name",
                  "result REAL NOT NULL - CGPA (2.50-4.00)",
                  "college TEXT NOT NULL - College/Department",
                  "board TEXT NOT NULL - Education Board",
                  "created_at TIMESTAMP - Record timestamp"]:
        p = tf.add_paragraph()
        p.text = field
        p.level = 0
        p.font.size = Pt(13)
    
    # Statistics
    p = tf.add_paragraph()
    p.text = "Student Data Statistics"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "Batch 2021: 75 students (IDs 2102002-2102080) | Avg CGPA: 3.24"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Batch 2020: 8 students | Batch 2019: 1 student"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Total: 84 students | Overall Avg CGPA: 3.25"
    p.level = 1
    p.font.size = Pt(13)
    p.font.bold = True
    
    # Universities
    p = tf.add_paragraph()
    p.text = "Universities: 24 public institutions (DU, BUET, PSTU, CUET, RUET, etc.)"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

def add_ui_slide():
    """Slide 6: User Interface"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "User Interface & Design Features"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PURPLE
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = "Modern, Stylish & Interactive UI"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    # Design elements
    p = tf.add_paragraph()
    p.text = "🎨 Visual Design Elements"
    p.font.size = Pt(18)
    p.font.color.rgb = BLUE
    
    for item in ["Glass Morphism: Translucent container with backdrop blur",
                 "Animated Background: Floating gradient orbs (purple, pink, blue)",
                 "6 Quick Command Buttons: Instant access to features",
                 "Status Indicator: Real-time connection status",
                 "Typing Animation: Processing indicator with bouncing dots",
                 "Message Bubbles: Different styles for user vs server"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(13)
    
    # UX features
    p = tf.add_paragraph()
    p.text = "💡 User Experience Features"
    p.font.size = Pt(18)
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "📱 Responsive | ⚡ Real-time | 🎯 One-click | ⌨️ Keyboard shortcuts | 😊 Emoji support"
    p.level = 1
    p.font.size = Pt(14)
    p.font.color.rgb = PINK

def add_implementation_slide():
    """Slide 7: Implementation"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Implementation & Technical Details"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PURPLE
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = "Chat Commands"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    # Commands
    commands = [
        ("search [ID]", "Search student by ID"),
        ("find [NAME]", "Search by student name"),
        ("list all", "View all 84 students"),
        ("university [NAME/CODE]", "Search university"),
        ("list universities", "View all 24 universities"),
        ("help", "Display all commands")
    ]
    
    for cmd, desc in commands:
        p = tf.add_paragraph()
        p.text = f"{cmd} → {desc}"
        p.level = 0
        p.font.size = Pt(14)
    
    # Technical details
    p = tf.add_paragraph()
    p.text = "Technical Stack"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "Backend: Flask REST API with regex command parser"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Frontend: JavaScript AJAX with async/await"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Database: SQLite3 with CRUD operations"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Features: CORS enabled, error handling, input validation, JSON responses"
    p.level = 1
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY

def add_conclusion_slide():
    """Slide 7: Deployment & Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Deployment & Conclusion"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PURPLE
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = "Deployment Status"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    # Deployment
    p = tf.add_paragraph()
    p.text = "✅ Git Repository: Version controlled and organized"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "✅ GitHub: https://github.com/safiullah-foragy/Socket_Programming_Client_Server_Chatting"
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "✅ Production Ready: Procfile, requirements.txt, runtime.txt"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "✅ Deployment Options: Render.com, Railway.app, Heroku"
    p.font.size = Pt(14)
    
    # Skills
    p = tf.add_paragraph()
    p.text = "Technologies Mastered"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "Backend: Python Flask, SQLite | Frontend: HTML5, CSS3, JavaScript"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Concepts: Client-Server, REST API, HTTP, JSON, MVC pattern"
    p.level = 1
    p.font.size = Pt(13)
    
    # Future
    p = tf.add_paragraph()
    p.text = "Future Enhancements"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "🔐 User authentication | 📊 Admin dashboard | 📱 Mobile app | 🤖 AI chatbot"
    p.level = 1
    p.font.size = Pt(14)
    p.font.color.rgb = PINK
    
    # Thank you
    p = tf.add_paragraph()
    p.text = "Thank You! Questions?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

# Generate all slides
print("🎨 Creating PowerPoint presentation with Socket Programming focus...")
add_title_slide()
print("✅ Slide 1: Title slide created")
add_overview_slide()
print("✅ Slide 2: Overview slide created")
add_socket_programming_slide()
print("✅ Slide 3: Socket Programming Components slide created")
add_architecture_slide()
print("✅ Slide 4: Architecture slide created")
add_database_slide()
print("✅ Slide 5: Database slide created")
add_ui_slide()
print("✅ Slide 6: UI slide created")
add_implementation_slide()
print("✅ Slide 7: Implementation slide created")
add_conclusion_slide()
print("✅ Slide 8: Conclusion slide created - BONUS SLIDE!")

# Save presentation
output_file = "Student_Result_University_System_Presentation.pptx"
prs.save(output_file)
print(f"\n🎉 PowerPoint presentation created successfully!")
print(f"📁 File: {output_file}")
print(f"📊 Total slides: 8 (7 main + 1 bonus conclusion)")
print(f"🔌 Special focus: Socket Programming Components (Slide 3)")
print(f"\n✨ Your presentation is ready to use!")
